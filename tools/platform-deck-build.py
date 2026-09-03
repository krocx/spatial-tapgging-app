# deck_build.py — AR Operations Platform pitch deck (native shapes, python-pptx)
# Usage: python3 deck_build.py <out.pptx>   (design units: 1600x900 → 13.333x7.5 in)
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SX = 13.333/1600; SY = 7.5/900
def X(v): return Inches(v*SX)
def Y(v): return Inches(v*SY)
def rgb(h): h=h.lstrip('#'); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
S = None
BG = '#0f1633'

def rect(x,y,w,h,fill,line=None,lw=1.5,radius=0.12,shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = S.shapes.add_shape(shape, X(x),Y(y),X(w),Y(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    if line: sh.line.color.rgb = rgb(line); sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE: sh.adjustments[0] = radius
    sh.shadow.inherit = False; sh.text_frame.text = ''
    return sh

def text(x,y,w,h,runs,size=12,color='#cbd5e1',bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,spacing=None):
    tb = S.shapes.add_textbox(X(x),Y(y),X(w),Y(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = Inches(0.02)
    paras = runs if isinstance(runs, list) else [(runs,size,color,bold)]
    for i,(t,sz,col,b) in enumerate(paras):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment = align
        if spacing: p.space_after = Pt(spacing)
        r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.color.rgb = rgb(col); r.font.bold = b; r.font.name = 'Calibri'
    return tb

def new():
    global S
    S = prs.slides.add_slide(prs.slide_layouts[6]); rect(0,0,1600,900,BG,shape=MSO_SHAPE.RECTANGLE); return S

def footer():
    text(60,868,900,22,'Proprietary & Confidential · Applied Materials · AR Operations Platform · v2026.4.45',8,'#475569')
    text(1140,868,400,22,'karthik_addagarla@amat.com',8,'#64748b',align=PP_ALIGN.RIGHT)

def badge(x,y,w,label,line,fg='#ffffff',size=7.5):
    b = rect(x,y,w,20,'#1e293b',line,1,0.5); b.text_frame.text = label
    p=b.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.runs[0]; r.font.size=Pt(size); r.font.bold=True; r.font.color.rgb=rgb(fg); r.font.name='Calibri'
    b.text_frame.margin_top=b.text_frame.margin_bottom=Inches(0)

# ── Platform map (approved layout) ───────────────────────────────────────────
def draw_map():
    new()
    text(60,28,900,40,'AR Operations Platform',24,'#ffffff',True)
    text(60,66,1000,26,'One spatial backbone · four products · authoring, content and governance built in',11,'#94a3b8')
    text(1140,32,400,22,'v2026.4.45 · September 2026',10,'#64748b',align=PP_ALIGN.RIGHT)
    text(1140,56,400,22,'Proprietary & Confidential · Applied Materials',9,'#475569',align=PP_ALIGN.RIGHT)
    def label(y,t): text(60,y-14,600,22,t,9,'#64748b',True)
    label(126,'PRODUCTS')
    products = [
      (60, '#1a1f4d','#818cf8','AR Work Instructions','AR OMS  ·  AR OJT (on-the-job training)','#a5b4fc',
       ['• Step-by-step AR guides with 3D ghost overlays','• Every step can be spatially validated','• Evidence, branches, resume & sign-off','• Kiosk shift start · Production # usage log']),
      (436,'#0f2a4a','#22d3ee','Spatial Validation','AR-guided · own models or 3rd-party via API','#67e8f9',
       ['• Author trains in AR from multiple angles','• Operator guided to the exact spot; live verdict','• Pass / fail states, regions of interest, evidence','• Pluggable engine — ours or any model via API']),
      (812,'#3a1f0f','#fb923c','GembaWalks','','#fdba74',
       ['• Location-tagged audit walks in AR','• Findings, defect categories, resolution','• Ghost re-localization at each checkpoint','• Walk reports & open-finding tracking']),
      (1188,'#3a0f14','#f87171','iLOTO','','#fca5a5',
       ['• Spatial Lockout / Tagout points & 3D locks','• Apply / remove checklists, live lock status','• Training quiz gate & certification','• AR LOTO energy-flow maps, EHS board']),
    ]
    for x,fill,line,title,sub,subcol,bul in products:
        rect(x,138,352,158,fill,line,2)
        text(x+16,146,330,30,title,15,'#ffffff',True)
        if sub: text(x+16,174,330,20,sub,8.5,subcol,True)
        text(x+16,196,330,100,[(b,9.5,'#cbd5e1',False) for b in bul],spacing=2)
    badge(1318,150,136,'EARLY PROTOTYPE','#f87171','#fca5a5',7)
    label(334,'AUTHORING & CONTENT')
    auth = [
      (60,470,'#14213d','#3b82f6','Roadmap & Procedure Designer','Visual canvas for roadmaps and procedures · role-typed nodes, branches & recovery paths · compile straight to an AR guide','Preview mode · round-trip edit of any guide · pre-flight checks','#93c5fd'),
      (554,546,'#0f2f2a','#2dd4bf','Content Library','Guide Library · 3D Model Library with CAD import (GLB → USDZ) · Guide import from Excel, JSON and MES sources · per-anchor model kits','Share guides per user · publish / draft control · ghost opacity per step','#5eead4'),
      (1124,416,'#1f1a3d','#a78bfa','Feature Catalogue','Docs-as-data map of every capability, with flows, API reference and architecture views','Ask SIB — docs-grounded assistant','#c4b5fd'),
    ]
    for x,w,fill,line,title,body,foot,fc in auth:
        rect(x,346,w,122,fill,line,1.5)
        text(x+16,354,w-30,26,title,13,'#ffffff',True)
        text(x+16,382,w-30,52,body,9.5,'#cbd5e1')
        text(x+16,436,w-30,24,foot,9.5,fc)
    label(506,'OPERATE & GOVERN')
    rect(60,518,1480,96,'#111a33','#334155',1.5)
    for x,w,t,d in [(84,270,'Web Portal','Tile home · every product on one screen'),
                    (370,290,'Completions & Usage Log','Step timing, verdicts, evidence · Excel'),
                    (676,290,'Access & Roles (UAM)','Owner → Technician · product entitlements'),
                    (982,270,'Admin & Ops','Ops log · backups · site lock-down'),
                    (1268,248,'Kiosk Mode','Shared iPads · employee ID sign-in')]:
        rect(x,538,w,56,'#1e293b',None,radius=0.18)
        text(x+10,540,w-16,24,t,10.5,'#ffffff',True); text(x+10,564,w-16,26,d,8,'#94a3b8')
    rect(60,640,1480,96,'#27317a','#6366f1',2)
    text(84,648,560,30,'SIB — Spatial Intelligence Backbone',14,'#ffffff',True)
    text(84,678,560,26,'The shared engine every product runs on · one server, on-prem or cloud',9.5,'#c7d2fe')
    for x,a,b2 in [(640,'Anchors & QR','World maps & re-localization'),(900,'Perception engine','Live session stream & AI hints'),(1180,'Encrypted training data','.tag / .sib spatial data formats')]:
        text(x,652,270,24,a,9.5,'#e0e7ff',True); text(x,678,300,24,b2,9.5,'#e0e7ff',True)
    label(778,'HOW WE GOT HERE')
    ln = S.shapes.add_connector(1, X(60),Y(820),X(1540),Y(820)); ln.line.color.rgb=rgb('#334155'); ln.line.width=Pt(1.5)
    ms = [(120,'#22d3ee','Spatial\nValidation'),(300,'#fb923c','GembaWalks\n(Loc-Tags)'),(480,'#818cf8','AR Work\nInstructions'),
          (660,'#2dd4bf','3D Models &\nCAD import'),(840,'#3b82f6','Roadmap &\nProcedure Designer'),(1020,'#f87171','iLOTO\n& certification'),
          (1200,'#a78bfa','Feature Catalogue,\n.tag, lock-down · Aug ’26'),(1400,'#22c55e','UAM, Kiosk, Usage Log,\nvalidated work steps · Sep ’26')]
    for cx,col,t in ms:
        d = S.shapes.add_shape(MSO_SHAPE.OVAL, X(cx-7),Y(813),X(14),Y(14)); d.fill.solid(); d.fill.fore_color.rgb=rgb(col); d.line.fill.background(); d.shadow.inherit=False
        text(cx-95,832,190,40,[(line,8.5,'#cbd5e1',False) for line in t.split('\n')],align=PP_ALIGN.CENTER)
    text(900,876,640,20,'Native iOS (ARKit) app + SIB server + web portal · Render (cloud) or on-prem Windows service',8,'#475569',align=PP_ALIGN.RIGHT)

# ── Deck ─────────────────────────────────────────────────────────────────────
def build_deck():
    # 1 Title
    new()
    text(80,250,1440,40,'APPLIED MATERIALS · AR OPERATIONS PLATFORM',11,'#818cf8',True)
    text(80,300,1440,200,[('Work Instructions that assist technicians as they do it —',34,'#ffffff',True),('and help to validate immediately.',34,'#22d3ee',True)])
    text(80,520,1200,90,'Four products on one spatial backbone — built in-house, runs on-prem or in the cloud, and proves every step with evidence.',15,'#cbd5e1')
    for i,(t,d) in enumerate([('Guided','Steps appear on the equipment itself, in the technician’s field of view.'),('Validated','Critical steps are checked spatially before moving on; verdict and score recorded.'),('Evidenced','Every run leaves a usage log — timings, photos, sign-off — exportable to Excel.')]):
        x=80+i*480; rect(x,640,440,140,'#1e293b','#334155',1); text(x+18,652,400,34,t,18,'#ffffff',True); text(x+18,690,404,80,d,11,'#94a3b8')
    footer()

    # 2 Why AR
    new()
    text(60,40,900,40,'WHY AR, WHY NOW',11,'#818cf8',True); text(60,72,1200,50,'The facts — and where Applied can be ahead',26,'#ffffff',True)
    facts=[('−40%','fewer procedural errors and 32% faster task completion with AR-guided instructions vs. paper or static methods.','PTC industrial AR benchmark 2025 (via OxMaint)'),
           ('~90%','improvement in first-time quality with AR instructions in Boeing wing assembly; wiring production time cut 25%.','Iowa State / Boeing — SME Advanced Manufacturing; AR Insider'),
           ('−60%','training time at Volvo Group; Fujitsu cut operator training from 3–5 days to about an hour.','PTC case studies: Volvo Group; Fujitsu'),
           ('67,000','U.S. semiconductor technician / engineer / CS roles at risk of going unfilled by 2030 — 39% of them technicians.','SIA / Oxford Economics, 2023'),
           ('60+','customers supported by ASML through AR remote assistance on lithography tools — AR is already normal in the fab, as a vendor’s support channel.','optics.org; Microsoft / ASML')]
    y=140
    for n,b,src in facts:
        text(60,y,150,50,n,24,'#ffffff',True); text(220,y,640,70,[(b,11.5,'#cbd5e1',False),(src,8.5,'#64748b',False)]); y+=122
    rect(920,140,620,610,'#1b2050','#3b3f8a',1.5)
    text(944,156,570,34,'Where Applied can be ahead',17,'#ffffff',True)
    vis=[('Own the backbone, not a vendor’s app.','Anchors, world maps, validation and usage data are Applied’s — deployable on-prem inside the fab network.'),
         ('Validate, don’t just display.','Most AR tools show steps. Ours checks them spatially and records the verdict — never blocking the line, always logging.'),
         ('Pluggable intelligence.','Validation runs on our engine today and can call any internal or third-party model via API.'),
         ('One platform, many BUs.','OMS, on-the-job training, EHS audits and lockout/tagout share authoring, access control and reporting.'),
         ('Authoring in minutes.','Guides from Excel, JSON or MES exports — or the visual Procedure Designer, with round-trip editing.')]
    y=200
    for h,d in vis:
        text(944,y,570,90,[(h,12,'#ffffff',True),(d,10.5,'#c7d2fe',False)]); y+=106
    footer()

    # 3–7 Products
    prods=[('AR Work Instructions','AR OMS · AR OJT (on-the-job training)','POC-READY','#818cf8',
      'Procedures shown in AR on the actual tool — each step pinned where the work is, with 3D ghosts of the correct result.',
      'Procedure knowledge lives in people and PDFs. We wanted it to live on the equipment.',
      ['Faster ramp — the guide teaches while they work (AR OJT)','Fewer skipped steps; critical ones validated on the spot','Audit-ready record per Production #','Shared kiosk iPads — employee-ID sign-in'],
      'Live on cloud and in-house servers. Best first use: a procedure you already have in Excel.'),
     ('Spatial Validation','AR-guided · own models or third-party via API','POC-READY','#22d3ee',
      'Confirms hands-on work was done right: the technician stands where the author stood; the platform compares the live view with “correct”, from several angles.',
      'Vision systems validate parts on a line — nothing validated a technician’s work at a tool.',
      ['Catch the missing guard before the tool restarts','Objective verdict per step, on record','A failed check can be overridden — and it is logged','Pluggable engine — ours today, any model via API'],
      'Multi-angle training and guided capture live; scoring tuned with field tests.'),
     ('GembaWalks','Audit walks · findings · resolution','POC-READY','#fb923c',
      'Audit walks in AR — checkpoints tagged to real locations, findings with photos, every open finding tracked to resolution.',
      'Findings in spreadsheets die. Findings tied to a physical location stay honest and repeatable.',
      ['Same checkpoints, same order, every walk','Evidence and location on every finding','Ageing open findings visible at a glance'],
      'Live end to end. Ready for a team’s first real walk.'),
     ('iLOTO','Spatial Lockout / Tagout','EARLY PROTOTYPE','#f87171',
      'Lockout/Tagout points on the real equipment in AR — checklists, live lock status, energy-flow map, and a certification quiz gate.',
      'LOTO is the highest-consequence, most paper-bound procedure on the floor. We explored making it visual and verifiable.',
      ['Which breaker, which lock, which order — shown in place','EHS gets a live board and a certification gate'],
      'Early prototype: works end to end; not yet reviewed against EHS / regulatory requirements.'),
     ('SIB — Spatial Intelligence Backbone','The engine every product runs on','IN PRODUCTION USE','#6366f1',
      'QR anchors tie content to equipment, world maps let any iPad re-find it, plus validation, telemetry, access control and usage records — one server, on-prem or cloud.',
      'Every AR product needs the same 80%. Built once, each new product is a thin layer — GembaWalks and iLOTO took weeks.',
      ['Data can stay inside the fab network','One user list and role model across products','Authoring from Excel / JSON / MES or the visual designer'],
      'Running continuously; hardened for POCs. Spatial data formats patent-pending.')]
    for name,sub,bl,col,what,why,how,stand in prods:
        new()
        text(60,40,900,40,'THE PRODUCTS',11,col,True)
        text(60,72,1100,50,name,28,'#ffffff',True); text(60,120,900,26,sub,12,col,True)
        bc = {'POC-READY':'#d97706','EARLY PROTOTYPE':'#dc2626','IN PRODUCTION USE':'#16a34a'}[bl]
        b=rect(1300,80,240,30,'#1e293b',bc,1,0.5); b.text_frame.text=bl; p=b.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.runs[0]; r.font.size=Pt(10); r.font.bold=True; r.font.color.rgb=rgb('#ffffff'); r.font.name='Calibri'
        for i,(h,d) in enumerate([('WHAT IT IS',what),('WHY WE BUILT IT',why)]):
            x=60+i*750; rect(x,170,720,220,'#1e293b','#334155',1); text(x+18,182,680,26,h,10,col,True); text(x+18,212,684,170,d,12.5,'#cbd5e1')
        rect(60,410,720,300,'#1e293b','#334155',1); text(78,422,680,26,'HOW IT HELPS YOU',10,col,True); text(78,452,684,250,[('• '+h,12,'#cbd5e1',False) for h in how],spacing=4)
        rect(810,410,730,300,'#1e293b','#334155',1); text(828,422,680,26,'WHERE IT STANDS',10,col,True); text(828,452,694,160,stand,12.5,'#cbd5e1')
        rect(828,620,694,64,'#312e81','#6366f1',1); text(846,628,660,50,[('Reach out to AppliedX — Teams chat or karthik_addagarla@amat.com',12,'#ffffff',True),('Bring one procedure or audit; we can have it in AR in a working session.',9.5,'#c7d2fe',False)])
        footer()

    # 8 Map
    draw_map()

    # 9 Contact
    new()
    text(80,280,1440,40,'REACH OUT TO US TO LEARN MORE',11,'#818cf8',True)
    text(80,330,1440,80,'Bring one procedure. We’ll put it in AR together.',34,'#ffffff',True)
    text(80,430,1200,40,'AppliedX · AR Operations Platform · Applied Materials',16,'#cbd5e1')
    text(80,470,1200,50,'Teams chat  ·  karthik_addagarla@amat.com',20,'#22d3ee',True)
    text(80,560,1200,90,'Send a BU use case: product of interest · equipment / area · the problem you want solved.',13,'#94a3b8')
    footer()

if __name__ == '__main__':
    out = sys.argv[1] if len(sys.argv) > 1 else 'AR-Platform-Pitch.pptx'
    if '--map-only' in sys.argv: draw_map()
    else: build_deck()
    prs.save(out); print('saved', out, 'slides:', len(prs.slides))
